"""Document text extraction for Layer 0 pipeline.

Extracts text from PDF, DOCX, RTF, XLSX, and PPTX documents.

Engine priority per format:
    PDF  : pymupdf (fitz) > pdfplumber > PyPDF2
    DOCX : python-docx
    RTF  : striprtf
    XLSX : openpyxl
    PPTX : python-pptx

ALL imports are optional.  When no extraction library is installed for a
given format, ``extract_text_from_document`` returns an empty
``DocResult`` with appropriate warnings so the rest of the pipeline
continues without error.
"""

from __future__ import annotations

import io
import logging
import os
import re
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional dependency probing
# ---------------------------------------------------------------------------

_HAS_PYMUPDF = False
try:
    import fitz  # type: ignore[import-untyped]  # pymupdf

    _HAS_PYMUPDF = True
except ImportError:
    pass

_HAS_PDFPLUMBER = False
try:
    import pdfplumber  # type: ignore[import-untyped]

    _HAS_PDFPLUMBER = True
except ImportError:
    pass

_HAS_PYPDF2 = False
try:
    import PyPDF2  # type: ignore[import-untyped]

    _HAS_PYPDF2 = True
except ImportError:
    pass

_HAS_DOCX = False
try:
    import docx  # type: ignore[import-untyped]  # python-docx

    _HAS_DOCX = True
except ImportError:
    pass

_HAS_STRIPRTF = False
try:
    from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]

    _HAS_STRIPRTF = True
except ImportError:
    pass

_HAS_OPENPYXL = False
try:
    import openpyxl  # type: ignore[import-untyped]

    _HAS_OPENPYXL = True
except ImportError:
    pass

_HAS_PPTX = False
try:
    from pptx import Presentation  # type: ignore[import-untyped]  # python-pptx

    _HAS_PPTX = True
except ImportError:
    pass

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

#: Maximum number of pages to process (default 100, env-configurable).
MAX_PAGE_COUNT: int = int(os.getenv("L0_MAX_DOC_PAGES", 100))

#: Maximum extracted text size in bytes (default 1 MB, env-configurable).
MAX_TEXT_BYTES: int = int(os.getenv("L0_MAX_DOC_TEXT_BYTES", 1 * 1024 * 1024))

#: Maximum raw document size in bytes (default 50 MB, env-configurable).
MAX_DOC_BYTES: int = int(os.getenv("L0_MAX_DOC_BYTES", 50 * 1024 * 1024))

#: Supported document types and their magic-byte signatures.
_DOC_SIGNATURES: dict[str, list[bytes]] = {
    "pdf":  [b"%PDF"],
    "rtf":  [b"{\\rtf"],
    "docx": [b"PK\x03\x04"],   # ZIP-based Office Open XML
    "xlsx": [b"PK\x03\x04"],   # ZIP-based Office Open XML
    "pptx": [b"PK\x03\x04"],   # ZIP-based Office Open XML
}

# ---------------------------------------------------------------------------
# Dataclass
# ---------------------------------------------------------------------------


@dataclass
class DocResult:
    """Result of a document text extraction attempt.

    Attributes:
        text:       Extracted text (empty string if extraction failed).
        metadata:   Document metadata (title, author, dates, etc.).
        page_count: Number of pages/sheets/slides detected.
        engine:     Name of the library that produced the result.
        warnings:   Non-fatal issues encountered during extraction.
    """

    text: str = ""
    metadata: dict = field(default_factory=dict)
    page_count: int = 0
    engine: str = "none"
    warnings: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# PDF JavaScript / action detection
# ---------------------------------------------------------------------------

# Operators that indicate JavaScript or dangerous actions in PDF streams.
# These are scanned at the raw-byte level (fast, no PDF parser needed).
#
#   /JS, /JavaScript   — JavaScript code/dictionary
#   /OpenAction         — auto-execute action when PDF opens
#   /AA                 — additional actions (page open/close, focus, etc.)
#   /Launch             — launch external application
#   /SubmitForm         — form submission (data exfiltration risk)
#   /ImportData         — import external data into form fields
#
# Each entry maps a PDF operator (as bytes) to the anomaly flag category.

_PDF_JS_INDICATORS: list[tuple[bytes, str]] = [
    (b"/JS",         "pdf_javascript"),
    (b"/JavaScript", "pdf_javascript"),
    (b"/OpenAction", "pdf_auto_action"),
    (b"/AA",         "pdf_auto_action"),
    (b"/Launch",     "pdf_external_action"),
    (b"/SubmitForm", "pdf_external_action"),
    (b"/ImportData", "pdf_external_action"),
]

# Regex used to reduce false positives on short operators like /JS and /AA.
# In real PDFs these operators appear as PDF name tokens — they are preceded
# by whitespace, a line start, or certain delimiters (<<, [, /) and
# followed by whitespace or a delimiter.  We compile lazily.
_PDF_TOKEN_BOUNDARY = rb'(?:^|[\s<\[/])'
_PDF_TOKEN_TRAIL    = rb'(?:[\s>/\]\r\n(]|$)'


def detect_pdf_javascript(pdf_bytes: bytes | bytearray) -> dict:
    """Scan raw PDF bytes for JavaScript and dangerous action operators.

    This performs a fast byte-level scan — no PDF structure parsing is
    required.  It looks for the PDF name-object operators that indicate
    the presence of JavaScript code or auto-execute actions.

    Parameters
    ----------
    pdf_bytes : bytes | bytearray
        Raw bytes of a PDF file (or any data — non-PDF input is handled
        gracefully).

    Returns
    -------
    dict
        ``has_javascript`` : bool
            True if any JS/action indicator was found.
        ``js_indicators`` : list[str]
            Names of the operators found (e.g. ``["/JS", "/OpenAction"]``).
        ``anomaly_flags`` : set[str]
            Anomaly flag strings for downstream pipeline consumption.
            Possible values: ``"pdf_javascript"``, ``"pdf_auto_action"``,
            ``"pdf_external_action"``.
    """
    found_operators: list[str] = []
    flags: set[str] = set()

    if not pdf_bytes:
        return {
            "has_javascript": False,
            "js_indicators": found_operators,
            "anomaly_flags": flags,
        }

    for operator_bytes, flag_name in _PDF_JS_INDICATORS:
        # Quick pre-check: is the operator even present?
        if operator_bytes not in pdf_bytes:
            continue
        # Contextual check: verify the operator appears as a proper PDF
        # name token (not as a substring of a longer word or a comment).
        pattern = _PDF_TOKEN_BOUNDARY + re.escape(operator_bytes) + _PDF_TOKEN_TRAIL
        if re.search(pattern, pdf_bytes):
            op_str = operator_bytes.decode("ascii")
            if op_str not in found_operators:
                found_operators.append(op_str)
            flags.add(flag_name)

    return {
        "has_javascript": len(found_operators) > 0,
        "js_indicators": found_operators,
        "anomaly_flags": flags,
    }


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def detect_doc_type(data: bytes) -> str | None:
    """Detect document type from magic bytes.

    Returns one of ``"pdf"``, ``"rtf"``, ``"pk_office"`` (for DOCX/XLSX/
    PPTX -- caller must disambiguate), or ``None``.
    """
    trimmed = data[:16]
    if trimmed.startswith(b"%PDF"):
        return "pdf"
    if trimmed.startswith(b"{\\rtf"):
        return "rtf"
    if trimmed.startswith(b"PK\x03\x04"):
        return "pk_office"  # could be DOCX, XLSX, or PPTX
    return None


def extract_text_from_document(
    doc_data: bytes,
    doc_type: str,
    *,
    max_pages: int | None = None,
    max_text_bytes: int | None = None,
    max_doc_bytes: int | None = None,
) -> DocResult:
    """Extract text from raw document bytes.

    Parameters
    ----------
    doc_data:
        Raw bytes of the document file.
    doc_type:
        Document format: ``"pdf"``, ``"docx"``, ``"rtf"``, ``"xlsx"``,
        or ``"pptx"``.  Case-insensitive.
    max_pages:
        Override maximum page count.  ``None`` uses ``MAX_PAGE_COUNT``.
    max_text_bytes:
        Override maximum extracted text size.  ``None`` uses
        ``MAX_TEXT_BYTES``.
    max_doc_bytes:
        Override maximum raw document size.  ``None`` uses
        ``MAX_DOC_BYTES``.

    Returns
    -------
    DocResult
        Always returns a result -- never raises.  If no library is
        available or the document cannot be processed, the ``text``
        field is empty and ``engine`` is ``"none"``.
    """
    dtype = doc_type.lower().strip()
    eff_max_pages = max_pages if max_pages is not None else MAX_PAGE_COUNT
    eff_max_text = max_text_bytes if max_text_bytes is not None else MAX_TEXT_BYTES
    eff_max_doc = max_doc_bytes if max_doc_bytes is not None else MAX_DOC_BYTES

    # --- Guard: raw document size limit --------------------------------------
    if len(doc_data) > eff_max_doc:
        return DocResult(
            warnings=[
                "Document exceeds size limit ({} bytes > {} bytes)".format(
                    len(doc_data), eff_max_doc
                )
            ],
        )

    # --- Dispatch to format-specific extractor -------------------------------
    _EXTRACTORS = {
        "pdf":  _extract_pdf,
        "docx": _extract_docx,
        "rtf":  _extract_rtf,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
    }

    extractor = _EXTRACTORS.get(dtype)
    if extractor is None:
        return DocResult(
            warnings=["Unsupported document type: {}".format(dtype)],
        )

    result = extractor(doc_data, eff_max_pages)

    # --- PDF JavaScript / action detection -----------------------------------
    if dtype == "pdf":
        js_result = detect_pdf_javascript(doc_data)
        if js_result["has_javascript"]:
            result.warnings.append(
                "PDF JavaScript/action indicators found: {}".format(
                    ", ".join(js_result["js_indicators"])
                )
            )
            # Store anomaly flags in metadata for downstream consumption
            if not hasattr(result, "metadata") or result.metadata is None:
                result.metadata = {}
            result.metadata["pdf_js_detection"] = js_result
            # Also add the anomaly flag names directly to warnings so they
            # can be picked up by sanitizer.py as anomaly flags.
            for flag in sorted(js_result["anomaly_flags"]):
                result.warnings.append("flag:{}".format(flag))

    # --- Post-extraction text size cap ---------------------------------------
    if len(result.text.encode("utf-8", errors="replace")) > eff_max_text:
        # Truncate to roughly eff_max_text bytes (character-safe)
        truncated = result.text.encode("utf-8", errors="replace")[:eff_max_text]
        result.text = truncated.decode("utf-8", errors="ignore")
        result.warnings.append(
            "Extracted text truncated to {} bytes".format(eff_max_text)
        )

    return result


# ---------------------------------------------------------------------------
# PDF extraction (pymupdf > pdfplumber > PyPDF2)
# ---------------------------------------------------------------------------


def _extract_pdf(data: bytes, max_pages: int) -> DocResult:
    """Extract text from PDF bytes using the best available library."""
    if _HAS_PYMUPDF:
        return _extract_pdf_pymupdf(data, max_pages)
    if _HAS_PDFPLUMBER:
        return _extract_pdf_pdfplumber(data, max_pages)
    if _HAS_PYPDF2:
        return _extract_pdf_pypdf2(data, max_pages)
    return DocResult(
        warnings=[
            "No PDF library installed "
            "(install pymupdf, pdfplumber, or PyPDF2)"
        ],
    )


def _extract_pdf_pymupdf(data: bytes, max_pages: int) -> DocResult:
    """Extract PDF text via pymupdf (fitz)."""
    warns: list[str] = []
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        total_pages = len(doc)
        pages_to_read = min(total_pages, max_pages)
        if total_pages > max_pages:
            warns.append(
                "PDF has {} pages, limited to {}".format(total_pages, max_pages)
            )

        text_parts: list[str] = []
        for i in range(pages_to_read):
            page = doc[i]
            text_parts.append(page.get_text())

        metadata = {}
        if doc.metadata:
            metadata = {
                k: v for k, v in doc.metadata.items() if v
            }

        doc.close()
        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata=metadata,
            page_count=total_pages,
            engine="pymupdf",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("pymupdf error: {}".format(exc))
        # Try next engine in chain
        if _HAS_PDFPLUMBER:
            return _extract_pdf_pdfplumber(data, max_pages)
        if _HAS_PYPDF2:
            return _extract_pdf_pypdf2(data, max_pages)
        return DocResult(engine="none", warnings=warns)


def _extract_pdf_pdfplumber(data: bytes, max_pages: int) -> DocResult:
    """Extract PDF text via pdfplumber."""
    warns: list[str] = []
    try:
        pdf = pdfplumber.open(io.BytesIO(data))
        total_pages = len(pdf.pages)
        pages_to_read = min(total_pages, max_pages)
        if total_pages > max_pages:
            warns.append(
                "PDF has {} pages, limited to {}".format(total_pages, max_pages)
            )

        text_parts: list[str] = []
        for i in range(pages_to_read):
            page_text = pdf.pages[i].extract_text()
            if page_text:
                text_parts.append(page_text)

        metadata = {}
        if pdf.metadata:
            metadata = {
                k: v for k, v in pdf.metadata.items() if v
            }

        pdf.close()
        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata=metadata,
            page_count=total_pages,
            engine="pdfplumber",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("pdfplumber error: {}".format(exc))
        if _HAS_PYPDF2:
            return _extract_pdf_pypdf2(data, max_pages)
        return DocResult(engine="none", warnings=warns)


def _extract_pdf_pypdf2(data: bytes, max_pages: int) -> DocResult:
    """Extract PDF text via PyPDF2."""
    warns: list[str] = []
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        total_pages = len(reader.pages)
        pages_to_read = min(total_pages, max_pages)
        if total_pages > max_pages:
            warns.append(
                "PDF has {} pages, limited to {}".format(total_pages, max_pages)
            )

        text_parts: list[str] = []
        for i in range(pages_to_read):
            page_text = reader.pages[i].extract_text()
            if page_text:
                text_parts.append(page_text)

        metadata = {}
        if reader.metadata:
            metadata = {
                k: str(v) for k, v in reader.metadata.items() if v
            }

        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata=metadata,
            page_count=total_pages,
            engine="PyPDF2",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("PyPDF2 error: {}".format(exc))
        return DocResult(engine="none", warnings=warns)


# ---------------------------------------------------------------------------
# DOCX extraction (python-docx)
# ---------------------------------------------------------------------------


def _extract_docx(data: bytes, max_pages: int) -> DocResult:
    """Extract text from DOCX bytes via python-docx."""
    if not _HAS_DOCX:
        return DocResult(
            warnings=["python-docx is not installed -- cannot parse DOCX"],
        )

    warns: list[str] = []
    try:
        document = docx.Document(io.BytesIO(data))

        # Extract paragraph text
        text_parts: list[str] = []
        for para in document.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Also extract table cell text
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text)

        # Extract core properties as metadata
        metadata = {}
        try:
            props = document.core_properties
            if props.title:
                metadata["title"] = props.title
            if props.author:
                metadata["author"] = props.author
            if props.subject:
                metadata["subject"] = props.subject
            if props.created:
                metadata["created"] = str(props.created)
            if props.modified:
                metadata["modified"] = str(props.modified)
        except Exception:
            pass  # metadata extraction is best-effort

        # DOCX does not have a fixed page count -- estimate from section count
        section_count = len(document.sections) if document.sections else 1

        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata=metadata,
            page_count=section_count,
            engine="python-docx",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("python-docx error: {}".format(exc))
        return DocResult(engine="none", warnings=warns)


# ---------------------------------------------------------------------------
# RTF extraction (striprtf)
# ---------------------------------------------------------------------------


def _extract_rtf(data: bytes, max_pages: int) -> DocResult:
    """Extract text from RTF bytes via striprtf."""
    if not _HAS_STRIPRTF:
        return DocResult(
            warnings=["striprtf is not installed -- cannot parse RTF"],
        )

    warns: list[str] = []
    try:
        # RTF is text-based; decode to string first
        rtf_str = data.decode("utf-8", errors="replace")
        text = rtf_to_text(rtf_str)
        return DocResult(
            text=text.strip(),
            metadata={},
            page_count=1,  # RTF does not expose page count
            engine="striprtf",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("striprtf error: {}".format(exc))
        return DocResult(engine="none", warnings=warns)


# ---------------------------------------------------------------------------
# XLSX extraction (openpyxl)
# ---------------------------------------------------------------------------


def _extract_xlsx(data: bytes, max_pages: int) -> DocResult:
    """Extract cell text from XLSX bytes via openpyxl."""
    if not _HAS_OPENPYXL:
        return DocResult(
            warnings=["openpyxl is not installed -- cannot parse XLSX"],
        )

    warns: list[str] = []
    try:
        # data_only=True reads cached values, not formulas (no macro exec)
        wb = openpyxl.load_workbook(
            io.BytesIO(data), read_only=True, data_only=True
        )
        sheet_count = len(wb.sheetnames)
        sheets_to_read = min(sheet_count, max_pages)
        if sheet_count > max_pages:
            warns.append(
                "XLSX has {} sheets, limited to {}".format(
                    sheet_count, max_pages
                )
            )

        text_parts: list[str] = []
        for sheet_name in wb.sheetnames[:sheets_to_read]:
            ws = wb[sheet_name]
            text_parts.append("[Sheet: {}]".format(sheet_name))
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    text_parts.append("\t".join(cells))

        wb.close()
        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata={"sheet_names": wb.sheetnames},
            page_count=sheet_count,
            engine="openpyxl",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("openpyxl error: {}".format(exc))
        return DocResult(engine="none", warnings=warns)


# ---------------------------------------------------------------------------
# PPTX extraction (python-pptx)
# ---------------------------------------------------------------------------


def _extract_pptx(data: bytes, max_pages: int) -> DocResult:
    """Extract slide text from PPTX bytes via python-pptx."""
    if not _HAS_PPTX:
        return DocResult(
            warnings=["python-pptx is not installed -- cannot parse PPTX"],
        )

    warns: list[str] = []
    try:
        prs = Presentation(io.BytesIO(data))
        total_slides = len(prs.slides)
        slides_to_read = min(total_slides, max_pages)
        if total_slides > max_pages:
            warns.append(
                "PPTX has {} slides, limited to {}".format(
                    total_slides, max_pages
                )
            )

        text_parts: list[str] = []
        for i, slide in enumerate(prs.slides):
            if i >= slides_to_read:
                break
            text_parts.append("[Slide {}]".format(i + 1))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_parts.append(paragraph.text)

        return DocResult(
            text="\n".join(text_parts).strip(),
            metadata={},
            page_count=total_slides,
            engine="python-pptx",
            warnings=warns,
        )
    except Exception as exc:
        warns.append("python-pptx error: {}".format(exc))
        return DocResult(engine="none", warnings=warns)
